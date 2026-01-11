from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from wordfreq import zipf_frequency
from english_words import get_english_words_set
import random
import nltk

WIDTH = Inches(13.33)
HEIGHT = Inches(7.5)

ALL_NOUNS = [
    w.lower()
    for w, t in nltk.corpus.brown.tagged_words(tagset="universal")
    if t == "NOUN" and w.isalpha() and 3 <= len(w) <= 12
]


def get_words(count=5, difficulty="easy"):
    if difficulty == "easy":
        pool = [w for w in ALL_NOUNS if zipf_frequency(w, "en") >= 5]

    elif difficulty == "medium":
        pool = [w for w in ALL_NOUNS if 4 <= zipf_frequency(w, "en") < 5]

    elif difficulty == "hard":
        pool = [w for w in ALL_NOUNS if zipf_frequency(w, "en") < 4]

    else:
        pool = sorted(ALL_NOUNS, key=lambda w: zipf_frequency(w, "en"))

    return random.sample(pool, count)


def create_game():
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT

    easy_words = get_words(20, "easy")
    medium_words = get_words(20, "medium")
    hard_words = get_words(10, "hard")

    words = easy_words + medium_words + hard_words

    print(words)
    for word in words:
        word = word.upper()
        # Scramble logic
        char_list = list(word)
        random.shuffle(char_list)
        scrambled = "".join(char_list)
        if scrambled == word:  # Ensure it's actually scrambled
            scrambled = word[1:] + word[0]

        # Slide 1: Scrambled Word
        slide_scrambled = prs.slides.add_slide(prs.slide_layouts[6])
        add_centered_text(slide_scrambled, scrambled, RGBColor(0, 0, 0))

        # Slide 2: Unscrambled Word
        slide_answer = prs.slides.add_slide(prs.slide_layouts[6])
        add_centered_text(slide_answer, word, RGBColor(0, 128, 0))

    prs.save("scrambled_game.pptx")


def add_centered_text(slide, text, color):
    txBox = slide.shapes.add_textbox(0, HEIGHT / 2 - Inches(1), WIDTH, Inches(2))
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(100)
    p.font.bold = True
    p.font.color.rgb = color


if __name__ == "__main__":
    create_game()
